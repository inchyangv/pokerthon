"""Pokerthon 소개 프레젠테이션 생성 스크립트."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──────────────────────────────────────
BG = RGBColor(0x0F, 0x0F, 0x0F)
SURFACE = RGBColor(0x1A, 0x1A, 0x1A)
BORDER = RGBColor(0x2A, 0x2A, 0x2A)
WHITE = RGBColor(0xF0, 0xF0, 0xF0)
DIM = RGBColor(0x99, 0x99, 0x99)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
GOLD = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
DARK_GREEN_BG = RGBColor(0x14, 0x2B, 0x14)
DARK_GOLD_BG = RGBColor(0x2B, 0x24, 0x10)
DARK_RED_BG = RGBColor(0x2B, 0x14, 0x14)
DARK_BLUE_BG = RGBColor(0x14, 0x1E, 0x2B)

FONT = "맑은 고딕"
FONT_EN = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# ── Helper functions ────────────────────────────

def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, color=WHITE,
             bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.line_spacing = Pt(size * line_spacing)
    return txBox


def add_multiline(slide, left, top, width, height, lines, size=18, color=DIM,
                  alignment=PP_ALIGN.LEFT, line_spacing=1.5, font_name=FONT):
    """lines: list of (text, color, bold, size_override)"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, clr, bld, sz = item, color, False, size
        else:
            txt = item[0]
            clr = item[1] if len(item) > 1 else color
            bld = item[2] if len(item) > 2 else False
            sz = item[3] if len(item) > 3 else size
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(4)
        p.line_spacing = Pt(sz * line_spacing)
    return txBox


def add_rect(slide, left, top, width, height, fill_color=SURFACE, border_color=BORDER):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    shape.rotation = 0
    shape.adjustments[0] = 0.05
    return shape


def add_accent_bar(slide, left=1.0, top=1.0, width=0.6, height=0.06):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, body_lines,
             title_color=GREEN, fill=SURFACE, border=BORDER):
    add_rect(slide, left, top, width, height, fill, border)
    add_text(slide, left + 0.3, top + 0.25, width - 0.6, 0.5, title,
             size=20, color=title_color, bold=True)
    add_multiline(slide, left + 0.3, top + 0.8, width - 0.6, height - 1.0,
                  body_lines, size=16, color=DIM, line_spacing=1.4)


# ════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, 0, 2.5, 16, 1.5, "POKERTHON", size=72, color=GREEN,
         bold=True, alignment=PP_ALIGN.CENTER, font_name="Segoe UI")
add_text(slide, 0, 4.2, 16, 0.7, "AI vs AI  |  No-Limit Texas Hold'em",
         size=24, color=DIM, alignment=PP_ALIGN.CENTER)

cards_text = "A♥   K♠   Q♦   J♣   T♥"
add_text(slide, 0, 5.3, 16, 0.7, cards_text, size=28, color=RGBColor(0x55, 0x55, 0x55),
         alignment=PP_ALIGN.CENTER, font_name="Segoe UI")

add_text(slide, 0, 7.5, 16, 0.5, "2026. 03. 21", size=16, color=DIM,
         alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════
# SLIDE 2 — About
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, 1.0, 0.8, 10, 0.8, "[이름]", size=36, color=WHITE, bold=True)
add_text(slide, 1.0, 1.5, 10, 0.6, "[회사명]", size=22, color=DIM)

add_card(slide, 1.0, 2.5, 6.5, 3.0,
         "회사 소개",
         [
             "[회사에 대한 간단한 소개]",
             "",
             "[하는 일 / 도메인 / 기술 스택 등]",
         ])

add_card(slide, 8.0, 2.5, 6.5, 3.0,
         "왜 만들었나",
         [
             "→  AI 에이전트 개발 실력 겨루기",
             "→  전략 + 코딩",
             "→  오프라인 모임 이벤트",
         ])


# ════════════════════════════════════════════════
# SLIDE 3 — What is Pokerthon
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_accent_bar(slide)
add_text(slide, 1.0, 1.2, 14, 0.8, "Pokerthon이 뭔가요?", size=36, color=WHITE, bold=True)
add_multiline(slide, 1.0, 2.1, 14, 1.0, [
    ("각자 AI 포커 에이전트를 만들어서", WHITE, False, 24),
    ("서버에 접속시키고, 자동으로 포커를 플레이하게 하는 대회", WHITE, False, 24),
], line_spacing=1.5)

w = 4.3
gap = 0.35
start_x = 1.0
y = 3.5
h = 3.5

add_card(slide, start_x, y, w, h,
         "AI가 플레이",
         ["사람이 직접 치는 게 아닙니다.", "", "각자 만든 봇이", "대신 플레이합니다."],
         title_color=GREEN)

add_card(slide, start_x + w + gap, y, w, h,
         "REST API",
         ["언어 제한 없음.", "", "HTTP 요청만 보낼 수 있으면", "어떤 언어든 OK."],
         title_color=BLUE)

add_card(slide, start_x + (w + gap) * 2, y, w, h,
         "전략 대결",
         ["코딩 실력 + 포커 전략.", "", "최종 칩 보유량으로", "순위를 결정합니다."],
         title_color=GOLD)


# ════════════════════════════════════════════════
# SLIDE 4 — How It Works (4 steps)
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_accent_bar(slide)
add_text(slide, 1.0, 1.2, 14, 0.8, "어떻게 진행되나요?", size=36, color=WHITE, bold=True)

steps = [
    ("1", GREEN, "크레덴셜 수령", "운영진이 API Key +\nSecret Key를 발급합니다.\n이게 봇의 신분증입니다."),
    ("2", GOLD, "에이전트 개발", "API 문서를 보고 봇을 만드세요.\n게임 상태 확인 → 액션 제출.\n어떤 언어든 OK."),
    ("3", BLUE, "테이블에 착석", "봇을 서버에 연결하고\n테이블에 앉히면\n자동으로 게임이 시작됩니다."),
    ("4", RED, "대결", "봇들이 자동으로\n홀덤을 플레이합니다.\n최종 칩으로 순위를 매깁니다."),
]

w = 3.2
gap = 0.4
start_x = 1.0
y = 2.5

for i, (num, color, title, desc) in enumerate(steps):
    x = start_x + i * (w + gap)
    add_rect(slide, x, y, w, 4.5, SURFACE, BORDER)
    add_text(slide, x, y + 0.3, w, 0.9, num, size=52, color=color,
             bold=True, alignment=PP_ALIGN.CENTER, font_name="Segoe UI")
    add_text(slide, x + 0.2, y + 1.5, w - 0.4, 0.5, title, size=22,
             color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + 0.3, y + 2.3, w - 0.6, 2.0, desc, size=15,
             color=DIM, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════
# SLIDE 5 — Game Rules
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_accent_bar(slide)
add_text(slide, 1.0, 1.2, 14, 0.8, "게임 규칙", size=36, color=WHITE, bold=True)

rules = [
    ("게임", "No-Limit Texas Hold'em"),
    ("블라인드", "SB 1 / BB 2  (고정)"),
    ("칩", "관리자가 계정에 직접 지급. 착석 시 차감 없음 — 칩은 항상 계정 자산"),
    ("테이블", "최대 9명, 최소 2명이면 게임 시작. 테이블 스택은 40칩으로 시작"),
    ("턴 타임아웃", "10분 — 초과 시 자동 FOLD"),
    ("최소 레이즈", "ceil(현재 베팅 × 1.5)  ← 커스텀 룰"),
    ("올인", "항상 가능, 사이드팟 자동 생성"),
    ("특이사항", "체크레이즈 OK / 콜레이즈 금지 / 스트래들 없음"),
]

y_start = 2.4
for i, (label, value) in enumerate(rules):
    y = y_start + i * 0.65
    if i % 2 == 0:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.0), Inches(y - 0.05), Inches(14.0), Inches(0.6),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x15, 0x15, 0x15)
        shape.line.fill.background()

    add_text(slide, 1.3, y, 3.0, 0.5, label, size=18, color=WHITE, bold=True)
    add_text(slide, 4.5, y, 10.0, 0.5, value, size=18, color=DIM)


# ════════════════════════════════════════════════
# SLIDE 6 — Bot Architecture
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_accent_bar(slide)
add_text(slide, 1.0, 1.2, 14, 0.8, "봇 만드는 법", size=36, color=WHITE, bold=True)

add_text(slide, 1.0, 2.3, 6, 0.5, "게임 루프", size=22, color=GREEN, bold=True)
loop_steps = [
    "→  게임 상태를 폴링 (Long-Polling 지원)",
    "→  내 턴인지 확인",
    "→  전략에 따라 액션 결정",
    "→  액션 제출 (FOLD / CHECK / CALL / RAISE / ALL_IN)",
    "→  반복",
]
add_multiline(slide, 1.0, 2.9, 6.5, 3.0, loop_steps, size=16, color=DIM, line_spacing=1.6)

add_text(slide, 1.0, 5.6, 6, 0.5, "인증", size=22, color=GREEN, bold=True)
auth_steps = [
    "→  HMAC-SHA256 서명",
    "→  헤더 4개 (X-API-KEY, X-TIMESTAMP, X-NONCE, X-SIGNATURE)",
    "→  Python / JS 예시 코드 제공",
]
add_multiline(slide, 1.0, 6.2, 6.5, 2.0, auth_steps, size=16, color=DIM, line_spacing=1.6)

# Right side - code block
add_rect(slide, 8.0, 2.3, 7.0, 5.5, RGBColor(0x12, 0x12, 0x12), BORDER)
code_lines = [
    ("# 최소 동작 봇", RGBColor(0x55, 0x55, 0x55), False, 14),
    ("", DIM, False, 10),
    ("while True:", RGBColor(0xC0, 0x84, 0xFC), True, 15),
    ("    state = poll_state(table_no)", DIM, False, 15),
    ("", DIM, False, 10),
    ("    if not state[\"legal_actions\"]:", RGBColor(0xC0, 0x84, 0xFC), False, 15),
    ("        continue  # 내 턴 아님", RGBColor(0x55, 0x55, 0x55), False, 15),
    ("", DIM, False, 10),
    ("    # 전략은 여기에", RGBColor(0x55, 0x55, 0x55), False, 14),
    ("    action = decide(state)", RGBColor(0x60, 0xA5, 0xFA), False, 15),
    ("", DIM, False, 10),
    ("    submit_action(", DIM, False, 15),
    ("        hand_id=state[\"hand_id\"],", DIM, False, 15),
    ("        action=action,", DIM, False, 15),
    ("    )", DIM, False, 15),
]
add_multiline(slide, 8.3, 2.5, 6.4, 5.0, code_lines, size=15, color=DIM,
              line_spacing=1.3, font_name="Courier New")


# ════════════════════════════════════════════════
# SLIDE 7 — What We Provide
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_accent_bar(slide)
add_text(slide, 1.0, 1.2, 14, 0.8, "제공되는 것들", size=36, color=WHITE, bold=True)

cards_data = [
    (1.0, "API 연동 명세서", GREEN,
     ["전체 엔드포인트 레퍼런스", "인증 방법 상세 설명", "Python·JS 코드 예시", "", "INTEGRATION.md"]),
    (8.0, "API Playground", BLUE,
     ["브라우저에서 바로 API 테스트", "서명 생성기", "인터랙티브 Explorer", "", "/playground"]),
    (1.0, "연습용 봇", GOLD,
     ["TAG(타이트), LAG(루즈),", "FISH(콜링스테이션) 봇 상주", "언제든 붙어서 테스트 가능", "", "항상 대기중"]),
    (8.0, "관전 페이지", RED,
     ["실시간 게임 관전", "리더보드", "핸드 히스토리 리플레이", "", "/viewer"]),
]

for i, (x, title, color, body) in enumerate(cards_data):
    y = 2.3 if i < 2 else 5.2
    add_card(slide, x, y, 6.5, 2.6, title, body, title_color=color)


# ════════════════════════════════════════════════
# SLIDE 8 — Core APIs
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_accent_bar(slide)
add_text(slide, 1.0, 1.2, 14, 0.8, "핵심 API", size=36, color=WHITE, bold=True)

apis = [
    ("GET",  GREEN, "/v1/private/me",                      "내 정보 확인 (인증 테스트용)"),
    ("POST", BLUE,  "/v1/private/tables/{no}/sit",          "테이블에 착석 (스택 40칩으로 시작)"),
    ("GET",  GREEN, "/v1/private/tables/{no}/state",        "게임 상태 + 내 홀카드 조회  ★"),
    ("POST", BLUE,  "/v1/private/tables/{no}/action",       "액션 제출  ★"),
    ("POST", BLUE,  "/v1/private/tables/{no}/stand",        "테이블에서 이석"),
    ("POST", GOLD,  "/v1/private/test/refill",              "[테스트 기간 한정] 지갑 40칩으로 보충 — 3/31 만료"),
]

y_start = 2.3
row_h = 0.83
for i, (method, color, path, desc) in enumerate(apis):
    y = y_start + i * row_h

    if i % 2 == 0:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.0), Inches(y - 0.08), Inches(14.0), Inches(row_h - 0.05),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x15, 0x15, 0x15)
        shape.line.fill.background()

    badge_bg = DARK_GREEN_BG if method == "GET" else (DARK_GOLD_BG if method == "POST" and color == GOLD else DARK_BLUE_BG)
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.3), Inches(y), Inches(0.9), Inches(0.5),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = badge_bg
    badge.line.fill.background()
    tf = badge.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = method
    p.font.size = Pt(13)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER

    add_text(slide, 2.5, y, 6.5, 0.5, path, size=16, color=WHITE, font_name="Courier New")
    add_text(slide, 9.5, y, 5.5, 0.5, desc, size=15, color=DIM)

add_multiline(slide, 1.0, 7.5, 14, 1.0, [
    ("★ 이 두 개가 핵심입니다. 상태 확인 → 액션 제출 → 반복.", WHITE, False, 16),
    ("+ 공개 API (테이블 목록, 핸드 이력, 리더보드)는 인증 없이 사용 가능", DIM, False, 14),
], line_spacing=1.6)


# ════════════════════════════════════════════════
# SLIDE 9 — Timeline
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_accent_bar(slide)
add_text(slide, 1.0, 1.2, 14, 0.8, "일정", size=36, color=WHITE, bold=True)

bar = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(1.5), Inches(3.7), Inches(13.0), Inches(0.04),
)
bar.fill.solid()
bar.fill.fore_color.rgb = BORDER
bar.line.fill.background()

timeline_items = [
    ("3/21 TODAY", GREEN,  "킥오프",       "규칙 설명\n크레덴셜 발급\n카톡방 입장",       1.5),
    ("3/22~3/25",  BLUE,   "개발 기간",    "API 연동\n인증 테스트\n기본 봇 구현",         5.3),
    ("3/26~3/30",  GOLD,   "전략 고도화",  "연습 봇 상대 테스트\n전략 개선\n엣지 케이스", 9.1),
    ("3/31",       RED,    "대회",         "전원 착석\n실시간 대결\n최종 순위 발표",       12.5),
]

for label, color, title, desc, x in timeline_items:
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x + 0.6), Inches(3.52), Inches(0.3), Inches(0.3),
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()

    add_text(slide, x, 2.7, 3.0, 0.4, label, size=12, color=DIM,
             bold=True, alignment=PP_ALIGN.CENTER, font_name="Segoe UI")
    add_text(slide, x, 4.1, 3.0, 0.5, title, size=20, color=WHITE,
             bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, 4.7, 3.0, 1.5, desc, size=14, color=DIM,
             alignment=PP_ALIGN.CENTER)

add_rect(slide, 3.5, 7.0, 9.0, 1.2, DARK_GOLD_BG, GOLD)
add_multiline(slide, 3.8, 7.15, 8.5, 1.0, [
    ("개발 기간: 3/21 ~ 3/30  |  대회: 3/31 (화)", GOLD, True, 22),
    ("테스트 기간 중 서버는 항상 열려 있습니다. 칩 리필: POST /v1/private/test/refill", DIM, False, 14),
], line_spacing=1.5, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════
# SLIDE 10 — Tips
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_accent_bar(slide)
add_text(slide, 1.0, 1.2, 14, 0.8, "전략 힌트", size=36, color=WHITE, bold=True)

add_card(slide, 1.0, 2.3, 6.5, 5.5,
         "시작하기",
         [
             "→  일단 FOLD만 하는 봇부터 만들어보세요",
             "→  인증 + 게임 루프가 되면 반은 성공",
             "→  제공된 예시 코드를 복붙해서 시작",
             "→  포커 룰 모르면 AI에게 물어보세요",
         ], title_color=GREEN)

add_card(slide, 8.0, 2.3, 6.5, 5.5,
         "이기고 싶다면",
         [
             "→  포지션(버튼 위치) 활용",
             "→  팟 오즈 계산",
             "→  상대 베팅 패턴 분석",
             "→  블러프 전략",
             "→  핸드 히스토리 API로 상대 분석",
         ], title_color=GOLD)


# ════════════════════════════════════════════════
# SLIDE 11 — Join (KakaoTalk)
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, 0, 1.2, 16, 0.8, "참가하기", size=42, color=WHITE,
         bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, 3.5, 2.5, 9.0, 5.0, DARK_GREEN_BG, GREEN)

qr = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6.75), Inches(2.9), Inches(2.5), Inches(2.5),
)
qr.fill.solid()
qr.fill.fore_color.rgb = WHITE
qr.line.fill.background()
tf = qr.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "카톡방\nQR코드"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

add_text(slide, 3.5, 5.6, 9.0, 0.6, "카카오톡 그룹채팅 입장",
         size=26, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, 3.5, 6.2, 9.0, 0.5, "질문, 공지, 크레덴셜 발급 모두 여기서",
         size=18, color=DIM, alignment=PP_ALIGN.CENTER)
add_text(slide, 3.5, 6.7, 9.0, 0.4, "[카톡방 링크]",
         size=16, color=GREEN, alignment=PP_ALIGN.CENTER)

add_text(slide, 0, 7.8, 16, 0.5,
         "입장하시면 바로 API Key + Secret Key를 발급해드립니다",
         size=16, color=DIM, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════
# SLIDE 12 — Resources
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_accent_bar(slide)
add_text(slide, 1.0, 1.2, 14, 0.8, "참고 자료", size=36, color=WHITE, bold=True)

add_card(slide, 1.0, 2.3, 6.5, 3.5,
         "문서",
         [
             "→  INTEGRATION.md — API 연동 명세서 (필독)",
             "→  인증 방법 + 전체 API 레퍼런스",
             "→  Python / JavaScript 코드 예시",
             "→  게임 룰 요약 + FAQ",
         ], title_color=GREEN)

add_card(slide, 8.0, 2.3, 6.5, 3.5,
         "웹",
         [
             "→  /playground — API 테스트",
             "→  /viewer — 관전 페이지",
             "→  /v1/public/* — 공개 API (인증 불필요)",
             "→  서버 주소는 카톡방에서 공지",
         ], title_color=BLUE)

add_rect(slide, 1.0, 6.5, 13.8, 1.2, SURFACE, BORDER)
add_text(slide, 1.0, 6.7, 13.8, 0.8,
         "질문은 카톡방으로",
         size=20, color=DIM, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════
# SLIDE 13 — End
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, 0, 3.5, 16, 1.4, "Q & A", size=64, color=WHITE,
         bold=True, alignment=PP_ALIGN.CENTER, font_name="Segoe UI")
add_text(slide, 0, 5.2, 16, 0.6, "서버 주소 및 크레덴셜은 카톡방에서 공유됩니다",
         size=18, color=DIM, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════
output_path = "pokerthon.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
